import re
import os
import time
import datetime
import tempfile
import asyncio
import glob
import base64
import hashlib
import ipaddress
import socket
import functools
from typing import List, Dict, Any
from contextlib import asynccontextmanager
from urllib.parse import urlparse
from collections import defaultdict

import httpx
import aiosqlite
import trafilatura
from bs4 import BeautifulSoup
from readability import Document as ReadabilityDoc
import justext
from rank_bm25 import BM25Okapi
from mcp.server.fastmcp import FastMCP
import docker
from playwright.async_api import async_playwright
from cachetools import TTLCache

# Документы и OCR
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
import pytesseract
from PIL import Image
import io

# ---- КОНФИГ -------------------------------------------------------------
SEARXNG_URL = "http://localhost:8888"
USER_AGENT = "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0 Safari/537.36"
HEADERS = {"User-Agent": USER_AGENT}
FORUM_PAGE_STEP = 20
CACHE_TTL = 86400  # 24 часа для кэша БД
MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 МБ лимит на загрузку страницы
MAX_INPUT_SIZE = 1_000_000  # 1 МБ лимит на входные параметры инструментов (в байтах)
MAX_DOCUMENT_SIZE = 50 * 1024 * 1024  # 50 МБ лимит на размер документов
SECCOMP_PROFILE_PATH = "/etc/docker/seccomp.json"  # Опциональный seccomp профиль
ALLOWED_DOCUMENT_DIRS = [os.path.expanduser("~/documents"), "/tmp"]  # Разрешённые директории для документов

# Глобальные кэши в RAM
bm25_cache = TTLCache(maxsize=100, ttl=3600)

# Rate limiting: IP -> {tool_name: [timestamps]}
rate_limits = defaultdict(lambda: defaultdict(list))
RATE_LIMIT_WINDOW = 60  # 1 минута
RATE_LIMIT_MAX_CALLS = 30  # Максимум 30 вызовов в минуту

# Ленивая инициализация ресурсов
_http_client = None
_db_conn = None
_playwright = None
_browser = None

try:
    docker_client = docker.from_env()
except Exception as e:
    print(f"ВНИМАНИЕ: Docker не запущен: {e}")

# ---- УПРАВЛЕНИЕ ЖИЗНЕННЫМ ЦИКЛОМ (LIFESPAN) ------------------------------
@asynccontextmanager
async def app_lifespan(app):
    """Graceful shutdown для освобождения ресурсов."""
    yield
    global _http_client, _db_conn, _playwright, _browser
    if _http_client:
        await _http_client.aclose()
    if _db_conn:
        await _db_conn.close()
    if _browser:
        await _browser.close()
    if _playwright:
        await _playwright.stop()

mcp = FastMCP("mega_agent_v8", host="0.0.0.0", port=8100, lifespan=app_lifespan)

# ---- RATE LIMITING DECORATOR ---------------------------------------------
def rate_limit(tool_name: str, max_calls: int = RATE_LIMIT_MAX_CALLS, window: int = RATE_LIMIT_WINDOW):
    """Декоратор для rate limiting по IP (для локального сервера используем 'local')."""
    def decorator(func):
        @functools.wraps(func)
        async def wrapper(*args, **kwargs):
            client_ip = "local"  # Для локального MCP сервера
            now = time.time()
            
            # Удаляем старые записи
            rate_limits[client_ip][tool_name] = [
                ts for ts in rate_limits[client_ip][tool_name] 
                if now - ts < window
            ]
            
            # Проверяем лимит
            if len(rate_limits[client_ip][tool_name]) >= max_calls:
                raise Exception(f"Rate limit exceeded for {tool_name}. Max {max_calls} calls per {window} seconds.")
            
            # Добавляем текущий вызов
            rate_limits[client_ip][tool_name].append(now)
            
            return await func(*args, **kwargs)
        return wrapper
    return decorator

# ---- ГЕТТЕРЫ РЕСУРСОВ ----------------------------------------------------
async def get_client() -> httpx.AsyncClient:
    global _http_client
    if _http_client is None:
        limits = httpx.Limits(max_connections=50, max_keepalive_connections=20)
        _http_client = httpx.AsyncClient(headers=HEADERS, follow_redirects=True, timeout=15, limits=limits)
    return _http_client

async def get_db() -> aiosqlite.Connection:
    global _db_conn
    if _db_conn is None:
        _db_conn = await aiosqlite.connect("mcp_cache.db")
        await _db_conn.execute("PRAGMA journal_mode=WAL;")
        await _db_conn.execute("PRAGMA synchronous=NORMAL;")
        await _db_conn.execute("CREATE TABLE IF NOT EXISTS page_cache (url TEXT PRIMARY KEY, content TEXT, ts REAL)")
        await _db_conn.execute("DELETE FROM page_cache WHERE ts < ?", (time.time() - CACHE_TTL,))
        await _db_conn.commit()
    return _db_conn

async def get_browser():
    global _playwright, _browser
    if _browser is None:
        _playwright = await async_playwright().start()
        _browser = await _playwright.chromium.launch(headless=True)
    return _browser

# ---- ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ---------------------------------------------
def _validate_url(url: str) -> bool:
    """Полная защита от SSRF: проверка схемы, hostname и ВСЕХ IP адресов (IPv4/IPv6)."""
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            return False
        hostname = parsed.hostname
        if not hostname:
            return False
        
        # Получаем ВСЕ IP адреса (IPv4 и IPv6)
        addr_infos = socket.getaddrinfo(hostname, None)
        for addr_info in addr_infos:
            ip_str = addr_info[4][0]
            addr = ipaddress.ip_address(ip_str)
            if addr.is_private or addr.is_loopback or addr.is_link_local or addr.is_multicast:
                return False
        return True
    except Exception:
        return False

def _validate_file_path(file_path: str) -> tuple[bool, str]:
    """Проверяет, что файл существует, не превышает лимит и находится в разрешённой директории."""
    if not os.path.exists(file_path):
        return False, f"Файл не найден: {file_path}"
    
    # Проверяем размер файла
    file_size = os.path.getsize(file_path)
    if file_size > MAX_DOCUMENT_SIZE:
        return False, f"Файл слишком большой: {file_size} байт (максимум {MAX_DOCUMENT_SIZE} байт)"
    
    # Проверяем, что файл находится в разрешённой директории
    abs_path = os.path.abspath(file_path)
    allowed = False
    for allowed_dir in ALLOWED_DOCUMENT_DIRS:
        if abs_path.startswith(os.path.abspath(allowed_dir)):
            allowed = True
            break
    
    if not allowed:
        return False, f"Файл находится вне разрешённых директорий: {ALLOWED_DOCUMENT_DIRS}"
    
    return True, ""

def _clean(text: str, limit: int = 0) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    return text[:limit] if limit > 0 else text

def _tokenize(text: str) -> list[str]:
    return re.findall(r"\w+", text.lower())

def _bm25_top_indices(chunks: list[str], query: str, top_k: int, cache_key: str = None) -> list[int]:
    if not query or len(chunks) <= top_k: 
        return list(range(min(len(chunks), top_k)))
    
    full_cache_key = None
    if cache_key:
        # Простое и эффективное хэширование: SHA256 от всех чанков
        chunks_hash = hashlib.sha256("|||".join(chunks).encode()).hexdigest()
        full_cache_key = f"{cache_key}_{chunks_hash}"
        
        if full_cache_key in bm25_cache:
            bm25 = bm25_cache[full_cache_key]
        else:
            tokenized = [_tokenize(c) for c in chunks]
            if not any(tokenized): 
                return list(range(min(len(chunks), top_k)))
            bm25 = BM25Okapi(tokenized)
            bm25_cache[full_cache_key] = bm25
    else:
        tokenized = [_tokenize(c) for c in chunks]
        if not any(tokenized): 
            return list(range(min(len(chunks), top_k)))
        bm25 = BM25Okapi(tokenized)

    scores = bm25.get_scores(_tokenize(query))
    ranked = sorted(range(len(chunks)), key=lambda i: scores[i], reverse=True)[:top_k]
    return sorted(ranked)

async def _fetch_html(url: str, use_cache: bool = True) -> str:
    if not _validate_url(url):
        raise ValueError("Недопустимый URL (SSRF защита: разрешены только http/https на публичные IP).")
        
    db = await get_db()
    if use_cache:
        async with db.execute("SELECT content FROM page_cache WHERE url=?", (url,)) as cursor:
            row = await cursor.fetchone()
            if row: return row[0]
        
    client = await get_client()
    async with client.stream("GET", url) as r:
        r.raise_for_status()
        content_length = int(r.headers.get("Content-Length", 0))
        if content_length > MAX_FILE_SIZE:
            raise ValueError(f"Размер страницы {content_length} байт превышает лимит в {MAX_FILE_SIZE} байт.")
        text = await r.aread()
        text = text.decode('utf-8', errors='ignore')
    
    if use_cache:
        await db.execute("REPLACE INTO page_cache (url, content, ts) VALUES (?, ?, ?)", (url, text, time.time()))
        await db.commit()
    return text

# ---- ИНСТРУМЕНТЫ ВЕБА ----------------------------------------------------
@mcp.tool()
@rate_limit("web_search")
async def web_search(query: str, language: str = "auto", num_results: int = 8) -> str:
    """Поиск в интернете (чтение сниппетов)."""
    if len(query.encode("utf-8")) > MAX_INPUT_SIZE:
        return f"Ошибка: Запрос слишком длинный (максимум {MAX_INPUT_SIZE} байт)."
    try:
        client = await get_client()
        r = await client.get(f"{SEARXNG_URL}/search", params={"q": query, "format": "json", "language": language})
        r.raise_for_status()
        results = r.json().get("results", [])[:num_results]
    except Exception as e:
        return f"Ошибка поиска: {e}"

    if not results: return "Ничего не найдено."
    out = [f"{i}. {res.get('title', '')}\n   URL: {res.get('url', '')}\n   {_clean(res.get('content', ''), 500)}" for i, res in enumerate(results, 1)]
    return "\n\n".join(out)

@mcp.tool()
@rate_limit("fetch_page")
async def fetch_page(url: str, max_chars: int = 15000) -> str:
    """Каскадный парсер статьи (Trafilatura -> Readability -> jusText -> BS4)."""
    try:
        html = await _fetch_html(url)
        text = trafilatura.extract(html, include_comments=False, include_tables=True)
        if not text or len(text) < 200:
            text = BeautifulSoup(ReadabilityDoc(html).summary(), "lxml").get_text("\n", strip=True)
        if not text or len(text) < 200:
            text = "\n".join([p.text for p in justext.justext(html, justext.get_stoplist("Russian")) if not p.is_boilerplate])
        if not text or len(text) < 200:
            soup = BeautifulSoup(html, "lxml")
            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]): tag.decompose()
            text = soup.get_text("\n", strip=True)
        return _clean(text)[:max_chars]
    except Exception as e:
        return f"Ошибка загрузки {url}: {e}"

@mcp.tool()
@rate_limit("search_and_read")
async def search_and_read(query: str, language: str = "auto", num_results: int = 3) -> str:
    """Комбо-инструмент: ищет 15 результатов, ранжирует по BM25, читает лучшие."""
    if len(query.encode("utf-8")) > MAX_INPUT_SIZE:
        return f"Ошибка: Запрос слишком длинный (максимум {MAX_INPUT_SIZE} байт)."
    try:
        client = await get_client()
        r = await client.get(f"{SEARXNG_URL}/search", params={"q": query, "format": "json", "language": language})
        r.raise_for_status()
        results = r.json().get("results", [])[:15]
    except Exception as e:
        return f"Ошибка поиска: {e}"

    if not results: return "Ничего не найдено."
    snippets = [res.get('title', '') + " " + res.get('content', '') for res in results]
    top_indices = await asyncio.to_thread(_bm25_top_indices, snippets, query, num_results)
    best_results = [results[i] for i in top_indices]

    tasks = [fetch_page(res.get("url", "")) for res in best_results]
    pages_content = await asyncio.gather(*tasks)

    chunks = []
    for res, content in zip(best_results, pages_content):
        url, title = res.get("url", ""), res.get("title", "")
        if any(f in url for f in ["reddit.com", "4pda.to", "habr.com"]):
             chunks.append(f"### {title}\nURL: {url}\n\n[Форум. Используй fetch_thread]")
        else:
             chunks.append(f"### {title}\nURL: {url}\n\n{content}")
    return "\n\n====================\n\n".join(chunks)

@mcp.tool()
@rate_limit("browse_url")
async def browse_url(url: str, max_chars: int = 15000) -> str:
    """Рендерит JS-heavy сайты через Playwright и извлекает текст (для SPA/React)."""
    if not _validate_url(url):
        return "Ошибка: Недопустимый URL (SSRF защита)."
    browser = await get_browser()
    page = None
    try:
        page = await browser.new_page()
        page.set_default_timeout(10000)  # 10 секунд на все операции
        await page.goto(url, wait_until="networkidle", timeout=30000)
        await page.evaluate("""() => {
            const tags = ['script', 'style', 'nav', 'footer', 'header', 'aside', 'noscript'];
            tags.forEach(tag => {
                document.querySelectorAll(tag).forEach(el => el.remove());
            });
        }""")
        text = await page.evaluate("() => document.body.innerText")
        return _clean(text)[:max_chars]
    except Exception as e:
        return f"Ошибка рендеринга {url}: {e}"
    finally:
        if page:
            try:
                await page.close()
            except:
                pass

# ---- ИНСТРУМЕНТЫ GITHUB --------------------------------------------------
@mcp.tool()
@rate_limit("fetch_github_repo")
async def fetch_github_repo(url: str) -> str:
    """Читает структуру файлов и содержимое README.md из GitHub репозитория."""
    if not _validate_url(url):
        return "Ошибка: Недопустимый URL."
    match = re.search(r"github\.com/([^/]+)/([^/]+)", url)
    if not match: return "Ошибка: Ожидается URL вида https://github.com/owner/repo"
    
    owner, repo = match.groups()
    repo = repo.replace(".git", "")
    client = await get_client()
    
    try:
        r_repo = await client.get(f"https://api.github.com/repos/{owner}/{repo}")
        r_repo.raise_for_status()
        default_branch = r_repo.json().get("default_branch", "main")
        
        r_tree = await client.get(f"https://api.github.com/repos/{owner}/{repo}/git/trees/{default_branch}?recursive=1")
        tree_items = r_tree.json().get("tree", []) if r_tree.status_code == 200 else []
        tree_paths = [item["path"] for item in tree_items if item["type"] == "blob"]
        tree_display = "\n".join(tree_paths[:50]) + ("\n... (обрезано)" if len(tree_paths) > 50 else "")

        readme_text = "README не найден."
        r_readme = await client.get(f"https://api.github.com/repos/{owner}/{repo}/readme")
        if r_readme.status_code == 200:
            content = r_readme.json().get("content", "")
            if content:
                readme_text = base64.b64decode(content).decode('utf-8', errors='ignore')

        return f"# Репозиторий: {owner}/{repo}\n\n## Структура файлов:\n{tree_display}\n\n## README.md:\n{readme_text[:10000]}"
    except Exception as e:
        return f"Ошибка обращения к GitHub API: {e}"

@mcp.tool()
@rate_limit("fetch_github_file")
async def fetch_github_file(url: str, file_path: str) -> str:
    """Читает конкретный файл из репозитория GitHub по его пути (например, src/main.py)."""
    if not _validate_url(url):
        return "Ошибка: Недопустимый URL."
    if len(file_path.encode("utf-8")) > 500:
        return "Ошибка: Путь к файлу слишком длинный."
    match = re.search(r"github\.com/([^/]+)/([^/]+)", url)
    if not match: return "Ошибка: Ожидается URL вида https://github.com/owner/repo"
    
    owner, repo = match.groups()
    repo = repo.replace(".git", "")
    client = await get_client()
    
    try:
        r_repo = await client.get(f"https://api.github.com/repos/{owner}/{repo}")
        r_repo.raise_for_status()
        default_branch = r_repo.json().get("default_branch", "main")
        
        api_url = f"https://api.github.com/repos/{owner}/{repo}/contents/{file_path}?ref={default_branch}"
        r_file = await client.get(api_url)
        r_file.raise_for_status()
        
        data = r_file.json()
        if data.get("encoding") == "base64" and data.get("content"):
            content = base64.b64decode(data["content"]).decode('utf-8', errors='ignore')
            return f"# Файл: {file_path}\n\n{content[:20000]}"
        else:
            return f"Не удалось декодировать файл {file_path} (возможно, он бинарный или слишком большой)."
    except httpx.HTTPStatusError as e:
        if e.response.status_code == 404:
            return f"Файл {file_path} не найден в репозитории."
        return f"Ошибка HTTP при чтении файла: {e}"
    except Exception as e:
        return f"Ошибка обращения к GitHub API: {e}"
# ---- ПАРСЕРЫ ФОРУМОВ -----------------------------------------------------
async def _fetch_reddit_chunks(url: str, sort: str, time_filter: str, max_comments: int) -> list[str]:
    if not _validate_url(url):
        raise ValueError("Недопустимый URL (SSRF защита).")
    json_url = url.split("?")[0].rstrip("/") + ".json"
    client = await get_client()
    r = await client.get(json_url, params={"sort": sort, "t": time_filter, "limit": 100})
    r.raise_for_status()
    data = r.json()
    post = data[0]["data"]["children"][0]["data"]
    chunks = [f"ЗАГОЛОВОК: {post.get('title')}\nПОСТ: {_clean(post.get('selftext', ''))}"]
    budget = [max_comments]

    def walk(children, depth=0):
        for c in children:
            if budget[0] <= 0 or depth > 4: return
            if c.get("kind") != "t1": continue
            d = c["data"]
            body = _clean(d.get("body", ""))
            if body and body not in ("[deleted]", "[removed]"):
                chunks.append(f"{'—' * depth} [{d.get('author', 'unknown')}] {body}")
                budget[0] -= 1
            replies = d.get("replies")
            if isinstance(replies, dict) and "data" in replies: walk(replies["data"]["children"], depth + 1)
    walk(data[1]["data"]["children"])
    return chunks

async def _fetch_4pda_thread(base_url: str, mode: str, max_pages: int) -> list[str]:
    if not _validate_url(base_url):
        raise ValueError("Недопустимый URL (SSRF защита).")
    first_html = await _fetch_html(base_url, use_cache=False)
    sts = [int(m) for m in re.findall(r"[&?]st=(\d+)", first_html)]
    max_st = max(sts) if sts else 0
    
    pages_to_fetch = [0]
    if max_st > 0:
        if mode == "recent":
            pages_to_fetch = list(range(max(0, max_st - max_pages * FORUM_PAGE_STEP), max_st + 1, FORUM_PAGE_STEP))
        elif mode == "full":
            pages_to_fetch = list(range(0, max_st + 1, FORUM_PAGE_STEP))[:max_pages]
        else: # first_last
            recent = max(FORUM_PAGE_STEP, max_st - (max_pages - 2) * FORUM_PAGE_STEP)
            pages_to_fetch += list(range(recent, max_st + 1, FORUM_PAGE_STEP))

    urls = [base_url if st == 0 else f"{base_url}&st={st}" for st in sorted(list(set(pages_to_fetch)))]
    
    async def fetch_one(u):
        try:
            html = await _fetch_html(u, use_cache=(u != urls[-1]))
            soup = BeautifulSoup(html, "lxml")
            return [_clean(b.get_text(" ", strip=True)) for i, b in enumerate(soup.select("div.post_body")) if not (u != urls[0] and i == 0)]
        except: return []

    results = await asyncio.gather(*[fetch_one(u) for u in urls])
    return [post for sublist in results for post in sublist if len(post) > 20]

@mcp.tool()
@rate_limit("fetch_thread")
async def fetch_thread(url: str, query: str = "", mode: str = "first_last", sort: str = "top", time_filter: str = "all", top_k: int = 25) -> str:
    """Парсит форумы. Reddit: sort, time_filter. 4PDA: mode."""
    if not _validate_url(url):
        return "Ошибка: Недопустимый URL (SSRF защита)."
    if len(query.encode("utf-8")) > MAX_INPUT_SIZE:
        return f"Ошибка: Запрос слишком длинный (максимум {MAX_INPUT_SIZE} байт)."
    try:
        if "reddit.com" in url:
            chunks, header = await _fetch_reddit_chunks(url, sort, time_filter, 150), "Reddit-тред"
        elif "4pda.to" in url or "4pda.ru" in url:
            chunks, header = await _fetch_4pda_thread(url, mode, 10), "Тема 4PDA"
        else:
            return "Поддерживается только Reddit и 4PDA."
    except Exception as e:
        return f"Ошибка треда: {e}"

    if query:
        idx = await asyncio.to_thread(_bm25_top_indices, chunks, query, top_k, cache_key=url)
        selected = [chunks[i] for i in idx]
        return f"# {header} (Топ {len(selected)} постов)\n\n" + "\n\n---\n\n".join(selected)
    return f"# {header} (Начало)\n\n" + "\n\n---\n\n".join(chunks[:top_k])

# ---- ПЕСОЧНИЦА, СИСТЕМА И ВИЗУАЛ -----------------------------------------
@mcp.tool()
@rate_limit("get_current_time", max_calls=60)
async def get_current_time() -> str:
    """Возвращает текущее системное время."""
    return f"Системное время: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

def _run_docker(cfg: dict, code: str) -> str:
    """Безопасная песочница с правильным timeout, ulimits и опциональным seccomp."""
    try:
        safe_code = base64.b64encode(code.encode('utf-8')).decode('utf-8')
        cmd = cfg["cmd"].replace("{B64_CODE}", safe_code)
        
        # Ограничения на размер файлов и процессы
        ulimits = [
            docker.types.Ulimit(name="fsize", soft=10_485_760, hard=10_485_760),  # 10MB
            docker.types.Ulimit(name="nproc", soft=32, hard=32),
            docker.types.Ulimit(name="nofile", soft=64, hard=64),
        ]
        
        security_opts = ["no-new-privileges:true"]
        if os.path.exists(SECCOMP_PROFILE_PATH):
            security_opts.append(f"seccomp={SECCOMP_PROFILE_PATH}")
        
        container = docker_client.containers.run(
            image=cfg["image"],
            command=["sh", "-c", cmd],
            detach=True,
            network_disabled=True,
            mem_limit="128m",
            memswap_limit="128m",
            nano_cpus=500_000_000,
            pids_limit=64,
            read_only=True,
            user="65534:65534",
            security_opt=security_opts,
            cap_drop=["ALL"],
            ulimits=ulimits,
            tmpfs={'/tmp': 'rw,exec,size=64m,uid=65534,mode=1777'},
            stdout=True,
            stderr=True
        )
        
        try:
            exit_info = container.wait(timeout=10)
            exit_code = exit_info.get("StatusCode", -1)
            output = container.logs().decode('utf-8', errors='ignore')
            container.remove(force=True)
            return f"Exit code: {exit_code}\n\nВывод:\n{output[:5000]}"
        except Exception as e:
            try:
                container.kill()
                container.remove(force=True)
            except:
                pass
            return f"Ошибка выполнения (timeout или сбой): {str(e)}"
    except Exception as e:
        return f"Ошибка Docker: {e}"

@mcp.tool()
@rate_limit("execute_code", max_calls=10)
async def execute_code(language: str, code: str) -> str:
    """Бронебойная песочница через Docker (non-root, seccomp, ulimits, isolation)."""
    if len(code.encode("utf-8")) > MAX_INPUT_SIZE:
        return f"Ошибка: Код слишком длинный (максимум {MAX_INPUT_SIZE} байт)."
    configs = {
        "python": {"image": "python:3.12-alpine", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/m.py && python /tmp/m.py"},
        "node": {"image": "node:22-alpine", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/m.js && node /tmp/m.js"},
        "php": {"image": "php:cli-alpine", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/m.php && php /tmp/m.php"},
        "c": {"image": "gcc:14-bookworm", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/m.c && gcc /tmp/m.c -o /tmp/a.out && /tmp/a.out"},
        "cpp": {"image": "gcc:14-bookworm", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/m.cpp && g++ /tmp/m.cpp -o /tmp/a.out && /tmp/a.out"},
        "java": {"image": "openjdk:21-alpine", "cmd": "printf '%s' '{B64_CODE}' | base64 -d > /tmp/Main.java && cd /tmp && javac Main.java && java Main"}
    }
    if language not in configs: return f"Язык '{language}' не поддерживается."
    return await asyncio.to_thread(_run_docker, configs[language], code)

@mcp.tool()
@rate_limit("render_html_css")
async def render_html_css(html_content: str) -> str:
    """Рендерит HTML, сохраняет скриншот и отдает локальный путь (использует Persistent Browser)."""
    if len(html_content.encode("utf-8")) > MAX_INPUT_SIZE:
        return f"Ошибка: HTML слишком большой (максимум {MAX_INPUT_SIZE} байт)."
    try:
        # Фоновая очистка старых скриншотов
        temp_dir = tempfile.gettempdir()
        now = time.time()
        for f in glob.glob(os.path.join(temp_dir, "mcp_render_*.jpg")):
            if os.path.isfile(f) and os.stat(f).st_mtime < now - 3600:
                try: os.remove(f)
                except: pass

        fd, temp_path = tempfile.mkstemp(suffix=".jpg", prefix="mcp_render_")
        os.close(fd)
        
        browser = await get_browser()
        page = None
        try:
            page = await browser.new_page(viewport={"width": 1280, "height": 720})
            page.set_default_timeout(10000)
            await page.set_content(html_content, wait_until="networkidle")
            await page.screenshot(path=temp_path, type="jpeg", quality=75)
        finally:
            if page:
                try:
                    await page.close()
                except:
                    pass
            
        return f"Скриншот сохранен по пути:\n{temp_path}"
    except Exception as e:
        return f"Ошибка рендеринга: {e}"

# ---- ПАРСЕРЫ ДОКУМЕНТОВ --------------------------------------------------
def _process_document(file_path: str, max_chars: int, ocr_lang: str) -> str:
    """Универсальный парсер документов с улучшенным PDF-извлечением."""
    valid, error_msg = _validate_file_path(file_path)
    if not valid:
        return error_msg
        
    ext = file_path.split('.')[-1].lower()
    text = ""

    try:
        if ext == "pdf":
            doc = fitz.open(file_path)
            for page in doc: 
                # Улучшенное извлечение: блоки с сортировкой по координатам (Y, затем X)
                blocks = page.get_text("blocks")
                if blocks:
                    sorted_blocks = sorted(blocks, key=lambda b: (b[1], b[0]))  # y0, x0
                    page_text = "\n".join([b[4] for b in sorted_blocks if b[6] == 0])  # b[6]==0 это текст, не картинка
                    if len(page_text.strip()) > 50:
                        text += page_text + "\n"
                    else:
                        # OCR fallback
                        images = page.get_images(full=True)
                        if images:
                            for img_index, img in enumerate(images):
                                xref = img[0]
                                base_image = doc.extract_image(xref)
                                image_bytes = base_image["image"]
                                pil_img = Image.open(io.BytesIO(image_bytes))
                                text += f"\n[OCR Изображения {img_index}]:\n" + pytesseract.image_to_string(pil_img, lang=ocr_lang) + "\n"
                        else:
                            pix = page.get_pixmap(dpi=150)
                            pil_img = Image.open(io.BytesIO(pix.tobytes("png")))
                            text += f"\n[OCR Страницы целиком]:\n" + pytesseract.image_to_string(pil_img, lang=ocr_lang) + "\n"
            doc.close()
        elif ext == "docx":
            text = "\n".join([p.text for p in docx.Document(file_path).paragraphs])
        elif ext == "pptx":
            for slide in pptx.Presentation(file_path).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext in ["xlsx", "xls", "ods", "csv"]:
            df_dict = {"CSV": pd.read_csv(file_path)} if ext == "csv" else pd.read_excel(file_path, sheet_name=None, engine='odf' if ext == 'ods' else None)
            text = "\n\n".join([f"### Лист: {name}\n{df.to_markdown(index=False)}" for name, df in df_dict.items()])
        
        clean = _clean(text)
        return clean[:max_chars] + ("...\n[Обрезано]" if len(clean) > max_chars else "")
    except Exception as e:
        return f"Ошибка чтения: {e}"

@mcp.tool()
@rate_limit("read_document")
async def read_document(file_path: str, max_chars: int = 30000, ocr_lang: str = "rus+eng") -> str:
    """Читает PDF (с умным OCR и сортировкой блоков), Word, PPTX, Excel."""
    return await asyncio.to_thread(_process_document, file_path, max_chars, ocr_lang)

# ---- ЗАПУСК --------------------------------------------------------------
if __name__ == "__main__":
    import uvicorn
    from starlette.middleware.cors import CORSMiddleware

    original_init = uvicorn.Config.__init__
    def patched_init(self, app, **kwargs):
        cors_app = CORSMiddleware(app, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"], expose_headers=["*"])
        original_init(self, cors_app, **kwargs)

    uvicorn.Config.__init__ = patched_init
    
    mcp.run(transport="streamable-http")
